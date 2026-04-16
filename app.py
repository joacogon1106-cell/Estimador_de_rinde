"""
Estimador de Rendimiento de Cultivos - Backend v4
"""
import http.server, socketserver, json, os, io, base64, struct, math, zipfile, tempfile, re

PORT = int(os.environ.get("PORT", 5050))

try:
    import numpy as np, tifffile
    from scipy import stats
    import matplotlib; matplotlib.use('Agg')
    import matplotlib.pyplot as plt, matplotlib.patheffects as pe
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                     Table, TableStyle, Image as RLImage,
                                     HRFlowable, PageBreak)
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    LIBS_OK = True; LIBS_ERR = ''
except ImportError as e:
    LIBS_OK = False; LIBS_ERR = str(e)

# Logo embebido en base64
LOGO_B64 = '/9j/4AAQSkZJRgABAQAAAQABAAD/4gHYSUNDX1BST0ZJTEUAAQEAAAHIAAAAAAQwAABtbnRyUkdCIFhZWiAH4AABAAEAAAAAAABhY3NwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAA9tYAAQAAAADTLQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAlkZXNjAAAA8AAAACRyWFlaAAABFAAAABRnWFlaAAABKAAAABRiWFlaAAABPAAAABR3dHB0AAABUAAAABRyVFJDAAABZAAAAChnVFJDAAABZAAAAChiVFJDAAABZAAAAChjcHJ0AAABjAAAADxtbHVjAAAAAAAAAAEAAAAMZW5VUwAAAAgAAAAcAHMAUgBHAEJYWVogAAAAAAAAb6IAADj1AAADkFhZWiAAAAAAAABimQAAt4UAABjaWFlaIAAAAAAAACSgAAAPhAAAts9YWVogAAAAAAAA9tYAAQAAAADTLXBhcmEAAAAAAAQAAAACZmYAAPKnAAANWQAAE9AAAApbAAAAAAAAAABtbHVjAAAAAAAAAAEAAAAMZW5VUwAAACAAAAAcAEcAbwBvAGcAbABlACAASQBuAGMALgAgADIAMAAxADb/2wBDAAUDBAQEAwUEBAQFBQUGBwwIBwcHBw8LCwkMEQ8SEhEPERETFhwXExQaFRERGCEYGh0dHx8fExciJCIeJBweHx7/2wBDAQUFBQcGBw4ICA4eFBEUHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh4eHh7/wAARCABsAI4DASIAAhEBAxEB/8QAHAABAQEBAAMBAQAAAAAAAAAAAAcGBQEDBAII/8QAPRAAAQIFAgQDBQYDCAMAAAAAAQIDAAQFBhEHIRITMUFRYXEUIjKBkQgVI1KCoRaiwiQzNEJDYnKxkrLB/8QAGgEBAAMBAQEAAAAAAAAAAAAAAAMEBQIBBv/EADARAAEEAQIEBAQGAwAAAAAAAAEAAgMRBBIhEzFBUQVhcaEUIoGxFUKRwdHwMlKi/9oADAMBAAIRAxEAPwD+y4QhBEhCEESEIQRIRir41EplvTaaVJS7tXrTmyJKW3KSenERnHpuY4TdP1buQB6dq8nbEsvcMMI43QPM77/qilJnMDyyMF7hzrp6nkP1ULpgDpaLPkqlCJXU7KuCmU5c7P6s1WVbb3W69s2P54+eQmNTqXIoqdJq1NvSlncBOA6R5EYOfmfSIznvYakiI9KPsDfsueOQac0+xVchGRsS/qRdKlSYS5T6q1nnSMxssEdeH8wH18RGui5DNHMzXGbCmY9rxbSkIQiVdJCEIIkIQgiQhCCJCEIIkYbVe65yjy8rQqCjnV+qq5Uskf6STsVn/wCfM9o3MS3TZAubUq47wmPfbk3fYJDPRKRsoj5f+xijmyPpsMZovNX2HMn+9SoZnHZjeZ/pWk06siQtKRLzqhN1d8cc3OubqUo7kAncJ/76mMfLarVVy4m5xymMptR6oGntzO/M4/zk58wcY6d8xTLqdWzbFVeQDxok3lJx1yEGIhNyrCPsyybqccYnA8D34+apP/UUM0uxQI8c6Q1pd61Wx9bNqCYmMBrNqBP6LXanN/xLqZbVnOkqkEpVPTbYOAsDOAfkkj9UNNmxbOqNx2eySmQdQmelG87IzjIH/lj9Ij0UZ1bmvUo9Ne6p6gIU1xdSSkE/1R9bpCvtHs8kklNIPOx264z/ACxGKMvH/NxK+lVX7rkbu19dVfRdjUqxWbhaTVaUv2C4ZX35aabPCVkdEqI/Y9vTaPdpXdzlz0d1iot+z1mnr5E8yRg8Q24gOwOD6EGNjEsuRH8K62UesMDlylfQZSaA2BcGAD9eD9/GL2Q0Y0onZsCQHfXYH1B9lPIOG4PHI8/5VThCEaasJCEIIkIQgiQhCCJCEIIvXNKUiWdWgZUlBI9cRBtJ7ouWh206in2bNViSVNuOOTLDnvcZxlOMHpgRfY/nq2qdX6fdVet6mXm3b65ObLjMvMY5UwlZyFDPfHD2PWMXxQvZNE9hI5jautHrt0VPJLg9pHn2/dbyU1btmYWqn3BI1GiuOAoWicYPAQdiMjf9owlEt5qp3siw5S4jOWs0fvVtLWFBY2/D4h33x9TjMaSqz9/02ULd3WtTLspWPeelUBSwnxxj+n5xxqHR6LVJw1zSysrpNaZSSulTZ+Id0jOdth4j0jPnfJO9rZNyDyI0uI6irp30KgeXPIDt/Y19itvrFbyl09u7qXOqp9VobSnGnEJyFtjcoP749SO8fDpdKyVJozl/3PW2lz9ZaC1vzCggNo7IT4nYdPAACOHX9Q6xW6K7ZL1svM3LOZlHmlKCWwCPjTnxHbp3yY5cxQbUtFmVF8VN+4as02ES9IlllSGR2Tjw9cZ8DEkmREcjjwjYDckkAO5bjvW1AWunPbxNbffYX/K289rBRFTBlqBS6pXXgcf2ZghJ+Z3/AGjF6kXJclZdoLlTtGYorLNTaXLvOuZUtRPw4wMbb/KNDITWpFSlAmgUKk2bSgPdXMIAWE9iRj+kRnF0+sVPU23qLUbtTcnLd9sf5OOVLhByRscZOMduo8Y4ypsiZmkk04gf4ho5+fzFcyPe8USd/Kh77q/whCPqVppCEIIkIQgiQhCCJCEIIkSLXuz0zK5e7peRM4JUBFRl0kgusg/ECNwRuM+GD2iuwUAoFKgCDsQe8VszFZlRGN39KjliErdJUNtmmPKlxP6W3o4taUhblGqCwSP9uDt88fqjwU0i9KoZKdlVWdfksribcQChD6x09c/XwKo6N+6SzCKga7ZD/sc0lXMMqHOXhXi0r/L/AMTt6dIzdx15up0NNM1CplQpNfk/8HVm5YjJHTjxjbP5cjuMGPmpWugBjlbQ/wCD6H8h7LOcCz5XCvsfTsV2FBd5uLtS62xSb2po4pCfT7vtGNxuOoOM7eoxgiOJZdTeoD79JkLSdqF9qmFoefmVcwIH58np167Z653xHzVS8KXclmy8xVqh7Fd1HWDJTaEHMykHIyQNj69xnuY9tWr1VvCq0ms2lSaoxcjcuWZ6Zl0ANLyMZCunjuex8hEDp4y4PY63bURRcQfL/dvfqFwXtJsHf3r+QupdlMZlWFzupt5vT07wFbdGkHeEZ7J2GAPPA9TGp0Fs1dGpz1wz8r7POVAfgMHqwxnIG++TsfQDzj5dONJRIzqa5djqJ6ocXMRL8XGhCvzLUfjV+3rFajYwMAmQZEraI5Amz6knr5dFbggJdrcK+/1SEIRuq6kIQgiQhCCJCBOBkxGZ7V6v1GrzMraNsifl5cnK1JW4tSQccRCfhB+cVcrMixQOIefKhZUUkzY61KzQjH6YXXU7ppU3NVWjfdipZ7lElRAWQMq91QBTjI6xjaxq/VJytu06zbf+80NEjmqStZcAO6glPRPmTET/ABLHZG2QnZ3LY2fouXZEbWhxPNWKETXTfVAXDWTQa1Tfu2p+9ywCeBZT1SQrdKuux6x+NR9TZqhXEm3KDSRUajhPHx8RAUoZCEpTuo4wYfieNweNq2uvO+1d0+Ij0a72VNj8PstPtlt5pDqD1StII+hicae33dNbuY0OuWt7AoMqdW7hbfABsMpUN8kgbGKVE+PkR5LNTOXmK+67jkbILC5ooFCDnMFFpoX+b2VGfriOg0220gNtIShA6JSMARiNVNQWrLaYYZkVTs/MJK0JJKW0JG3Eo48ewjIU7WKtyVVlmLsttMjKzGCHEJWhaUk/HhXxD0xFWTxHEx5OGTR67cvUqN2RFG7SVaIRitUb+l7Mp8spuWE5OzeeQ2VcKQkYypR8Nxt3jAz2r960ttl2q2pKyrcwnjYU6h1sLHiCTv2j3I8Uxsd5Y87jnQJr1XsmTGw0Vc4RLZzUa45HTv8Aiaft5mVmVzyGGWXCoJcaUjiDgPXfePVYmpdw3RctOkUW80xIOhQmZnhcUkKSkk8KvhHQDByYfimPrayzbqrY9V58THYb1Kq8Ilt96rO0y4VW9bVJ+9Z5tXLdUriKePuhKU7qI7naPntXVucVcTVDu6iClPurDaXE8SeBR+HjSrcAnuDA+K4ok4Zd1q6NX2vkhyow7TarUIQjRVhDuMGIFM2U/KVOfq+nV5SgS24sOMCa5TjeCcoJ6KSD04sD1i9vNpdaW0sEpWkpVg42PmIlFQ0Lt96ZU5J1aoSjRP8AdkJcx5AkZ+uYyfFMaScN0M1VfWiPQqrkxufVC/rS9WmN61u7bauCk1AJenpaSUWZhtPCV8SVJAIG3FkdR1jlfZjnqbL/AHxKTDrTM65ylIDhCSpABBAz4HqPOKhYlm0izpB2WpgdccfIU++6rK3COnTYAb7DxjMXVo5btZqbtQlZmZpjjyitxDISpsqPUhJ6Z8jiKow8xghmNOey7BPfz8lHwpRodzIWIrT8vUvtGyTtHcQ8kTbAcW1ukqQn8Q5HXYEE+UaTU+zafW709ro9zylMuFLSHFyzrnAVY2StKhulWBjbPSNXYGnNCs95c3Kl6bnlJ4PaH8ZSnuEgbDP1jnXrpNQrlqztW9tnZKbeIU6UKC0qIGM4V06djHP4fNwH62BznO1VdV6HuvOA/QbFkm6WT0+vS7KZfzNn3JNN1JLjnJ5gUFqbVw8QIWPiHiDuIt8YOw9L6FalRFTQ/MT86kENuPYAbzsSlI7+ZzG8jR8MhniiInO97b3Q7X1U+Ox7W09R7V297hlrylrSoDstJlxLYW+8lJ4lOHbdWQlI8cRPNWk3BLzsnJ3FcUnWZlttSkiXweQCRsogDc46eUW/UPTej3jNNTsxMTEnOto5fNZweNOcgKB64z1jPq0NtsyTTIqVRS8kqLjwKcuZxgYxgAeXjGRn+H5s75BzBO3zUK7Uqs8Ezy7qPX9l8GvtvqnLSpVwtzDaDT2UtuIWrBWlfDjh8SCOnh6RmJmV1E1OkqU3MSUuintZLc1gNoOfdKzuSTt0Aig6xG1U0ijW9clYnpNtSwtC5dAUSEJ4eJYwdt+oHWJLeNMsel0tEza92Ts/P8wDl4KQE9ySEjGIg8TY1k73X8tDUA4Akjy3XGQAHk9NrFqna8ySKbpRT6e2oqRLTMuylR6kJQoZ/aNPoylKdMaGEpAywSceJWrMcei2+9fWkVFk7inJtp5XC+XU45iwkqCCcjukgxtLWozFv2/J0aWdcdalUcCVuY4iMk749Y2caFzsr4kCmlgA+6txsJk4lbEKIaSTEtTdZas3V3G2ZhaplttTpA/E5mSMnuRn1j2faGmpGevWiy8g42/NtthDvKIUQVODgScd+px5xRL70voF1zxqK3H5CeUAHHWMEOY6cSTsT57R89laTW/blTbqbj8xUptpXEyp8AIbV+YJHU+ZjPPh2VwjigDSXXqvpd8u6gOPLpMVbXzVBTkJGeuN48whH0y0UhCEESEIQRIQhBEhCEESEIQRcyvW/RK82lusUuVnQj4S6jKk+h6iORLadWRLvoeatuS40HKeIFQz6EkGNVCIX48L3anMBPoFwY2E2QvCQEpCUgAAYAHaPMIRMu0hCEESEIQRf//Z'


def _get_html():
    p = os.path.join(os.path.dirname(os.path.abspath(__file__)), "app.html")
    return open(p, "r", encoding="utf-8").read() if os.path.exists(p) else "<h1>No se encontro app.html</h1>"

# SHAPEFILE

def fmt(n, dec=2):
    """Numero con punto separador de miles y coma decimal (estilo Argentina)."""
    try:
        n = float(n)
        # Formatear con coma decimal y punto de miles
        partes = f"{n:,.{dec}f}".split('.')
        entero = partes[0].replace(',', '.')
        decimal = partes[1] if len(partes) > 1 else '00'
        return f"{entero},{decimal}"
    except:
        return str(n)


def convertir_rinde(valor_kgha, unidad):
    """Convierte kg/ha a la unidad seleccionada por el usuario."""
    if unidad == 't/ha':
        return valor_kgha / 1000.0
    elif unidad == 'qq/ha':
        return valor_kgha / 100.0
    else:  # kg/ha
        return valor_kgha

def leer_dbf(data):
    n=struct.unpack('<I',data[4:8])[0]; hdr=struct.unpack('<H',data[8:10])[0]; rsz=struct.unpack('<H',data[10:12])[0]
    fields=[]; pos=32
    while data[pos]!=0x0D:
        fields.append((data[pos:pos+11].replace(b'\x00',b'').decode('latin1').strip(), data[pos+16])); pos+=32
    records=[]; rpos=hdr
    for _ in range(n):
        rec={}; fpos=1
        for name,flen in fields: rec[name]=data[rpos+fpos:rpos+fpos+flen].decode('latin1').strip(); fpos+=flen
        records.append(rec); rpos+=rsz
    return records

def _pip_simple(px, py, poly):
    n=len(poly); inside=False; j=n-1
    for i in range(n):
        xi,yi=poly[i]; xj,yj=poly[j]
        if ((yi>py)!=(yj>py)) and (px<(xj-xi)*(py-yi)/(yj-yi)+xi): inside=not inside
        j=i
    return inside

def leer_shp(data):
    """Lee poligonos del shapefile. Detecta exteriores y huecos por centroide."""
    polys=[]; pos=100
    while pos<len(data):
        clen=struct.unpack('>I',data[pos+4:pos+8])[0]*2
        if struct.unpack('<I',data[pos+8:pos+12])[0]==5:
            npts=struct.unpack('<I',data[pos+48:pos+52])[0]
            npar=struct.unpack('<I',data[pos+44:pos+48])[0]
            parts=[struct.unpack('<I',data[pos+52+i*4:pos+56+i*4])[0] for i in range(npar)]
            ps2=pos+52+npar*4
            all_pts=[(struct.unpack('<d',data[ps2+i*16:ps2+i*16+8])[0],
                      struct.unpack('<d',data[ps2+i*16+8:ps2+i*16+16])[0]) for i in range(npts)]
            if npar==1:
                polys.append(all_pts)
            else:
                partes=[]
                for pi in range(npar):
                    s=parts[pi]; e=parts[pi+1] if pi+1<npar else npts
                    p=all_pts[s:e]
                    if len(p)>=3: partes.append(p)
                exteriores=[]
                for i,p in enumerate(partes):
                    cx=sum(pt[0] for pt in p)/len(p); cy=sum(pt[1] for pt in p)/len(p)
                    es_int=any(i!=j and _pip_simple(cx,cy,partes[j]) for j in range(len(partes)))
                    if not es_int: exteriores.append(p)
                if not exteriores: exteriores=partes
                polys.append(exteriores[0] if len(exteriores)==1 else {'multipart':exteriores})
        else:
            polys.append(None)
        pos+=8+clen
    return polys
def pip(px,py,poly):
    if isinstance(poly, dict) and 'multipart' in poly:
        return any(pip(px,py,p) for p in poly['multipart'])
    n=len(poly); inside=False; j=n-1
    for i in range(n):
        xi,yi=poly[i]; xj,yj=poly[j]
        if ((yi>py)!=(yj>py)) and (px<(xj-xi)*(py-yi)/(yj-yi)+xi): inside=not inside
        j=i
    return inside

def area_ha(poly):
    if isinstance(poly, dict) and 'multipart' in poly:
        return sum(area_ha(p) for p in poly['multipart'])
    n=len(poly); a=0
    for i in range(n): j=(i+1)%n; a+=poly[i][0]*poly[j][1]-poly[j][0]*poly[i][1]
    a=abs(a)/2; cy=sum(p[1] for p in poly)/n
    return a*(111320*math.cos(math.radians(cy)))*111320/10000

def norm(s): return re.sub(r'\s+',' ',str(s).strip().lower())

def col_lote(records):
    if not records: return None
    for c in ['LOTE','LOTEPLANO','LOTE_PLANO','NOMBRE','NAME','ID']:
        if c in records[0]: return c
    return list(records[0].keys())[0]

def buscar_poly(nombre, records, polys, col):
    nb=norm(nombre); nbs=nb.replace(' ','')
    cols=[c for c in [col,'LOTEPLANO','LOTE','LOTE_PLANO'] if records and c in records[0]]
    for rec,p in zip(records,polys):
        if not p: continue
        for c in cols:
            shp=norm(rec.get(c,''))
            if shp==nb or shp.replace(' ','')==nbs: return p
    for rec,p in zip(records,polys):
        if not p: continue
        for c in cols:
            shp=norm(rec.get(c,''))
            if shp and (nb in shp or shp in nb): return p
    return None

# IMAGENES
def detectar_banda(nombre):
    """Detecta la banda en nombres de archivo de Copernicus EO Browser y otros formatos."""
    n = nombre.upper()
    if not (n.endswith('.TIF') or n.endswith('.TIFF')): return None
    if 'B8A' in n or 'B08A' in n: return None
    # Formato Copernicus EO Browser: *_B03_(Raw).tiff, *_B08_(Raw).tiff
    if re.search(r'_B03[_.(]|_B3[_.(]', n): return 'B3'
    if re.search(r'_B04[_.(]|_B4[_.(]', n): return 'B4'
    if re.search(r'_B08[_.(]|_B8[_.(]', n): return 'B8'
    # Otros formatos genericos
    if '_B03' in n or 'B03_' in n: return 'B3'
    if '_B04' in n or 'B04_' in n: return 'B4'
    if '_B08' in n or 'B08_' in n: return 'B8'
    return None

def leer_tiff(data):
    arr = tifffile.imread(io.BytesIO(data)).astype(np.float32)

    # Metodo 1: tifffile nativo (mas compatible con compresion Copernicus)
    try:
        tf = tifffile.TiffFile(io.BytesIO(data))
        page = tf.pages[0]
        tags = {t.code: t.value for t in page.tags.values()}
        if 33922 in tags and 33550 in tags:
            tp = tags[33922]
            ps = tags[33550]
            return arr, float(tp[3]), float(tp[4]), float(ps[0]), float(ps[1])
    except Exception:
        pass

    # Metodo 2: lectura binaria directa
    try:
        buf = io.BytesIO(data)
        magic = buf.read(2)
        end = '<' if magic == b'II' else '>'
        buf.read(2)
        ifd_offset = struct.unpack(end+'I', buf.read(4))[0]
        buf.seek(ifd_offset)
        n_entries = struct.unpack(end+'H', buf.read(2))[0]
        tags = {}
        for _ in range(n_entries):
            tag_id = struct.unpack(end+'H', buf.read(2))[0]
            struct.unpack(end+'H', buf.read(2))
            count  = struct.unpack(end+'I', buf.read(4))[0]
            val_buf = buf.read(4)
            if tag_id in (33922, 33550):
                offset = struct.unpack(end+'I', val_buf)[0]
                pos = buf.tell()
                buf.seek(offset)
                vals = [struct.unpack(end+'d', buf.read(8))[0] for _ in range(count)]
                buf.seek(pos)
                tags[tag_id] = vals
        if 33922 in tags and 33550 in tags:
            return arr, tags[33922][3], tags[33922][4], tags[33550][0], tags[33550][1]
    except Exception:
        pass

    raise Exception(
        "No se pudo leer la georreferencia del TIFF. "
        "Descarga las bandas como 'Raw' desde Copernicus EO Browser.")

def leer_zip(zip_bytes):
    """Lee bandas del ZIP de Copernicus extrayendo a disco temporal."""
    import shutil
    b3=b4=b8=None; geo=None; log=[]; tiffs_originales=[]
    tmp_dir = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(io.BytesIO(zip_bytes)) as zf:
            tiffs_originales = [n for n in zf.namelist() if n.upper().endswith(('.TIF','.TIFF'))]
            for i, name in enumerate(tiffs_originales):
                base = os.path.basename(name)
                banda = detectar_banda(base)
                if not banda:
                    log.append(f"SKIP={base}")
                    continue
                nombre_limpio = f"banda_{i}_{banda}.tiff"
                ruta_tmp = os.path.join(tmp_dir, nombre_limpio)
                try:
                    with zf.open(name) as src, open(ruta_tmp, 'wb') as dst:
                        dst.write(src.read())
                    with open(ruta_tmp, 'rb') as f:
                        data = f.read()
                    arr, olon, olat, plon, plat = leer_tiff(data)
                    g = (olon, olat, plon, plat)
                    log.append(f"{banda}={base} OK geo=({olon:.4f},{olat:.4f})")
                    if banda=='B3' and b3 is None:
                        b3=arr; geo=g
                    elif banda=='B4' and b4 is None:
                        b4=arr
                        if geo is None: geo=g
                    elif banda=='B8' and b8 is None:
                        b8=arr
                        if geo is None: geo=g
                except Exception as ex:
                    log.append(f"{banda}={base} ERR={str(ex)[:200]}")
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    print(f"leer_zip resultado: {log}")
    print(f"b3={b3 is not None} b4={b4 is not None} b8={b8 is not None} geo={geo}")

    if b8 is None or geo is None:
        raise Exception(
            f"No se pudo leer las bandas del ZIP.\n"
            f"TIFFs en ZIP: {[os.path.basename(t) for t in tiffs_originales]}\n"
            f"Log detallado: {log}\n"
            f"b8={'OK' if b8 is not None else 'FALLO'}, geo={'OK' if geo is not None else 'FALLO'}")
    if b3 is None: b3=np.zeros_like(b8)
    if b4 is None: b4=np.zeros_like(b8)
    return b3,b4,b8,geo

def calc_idx(b3,b4,b8,indice):
    r={}
    with np.errstate(invalid='ignore',divide='ignore'):
        if indice in ('ndvi','ambos'): r['NDVI']=np.where((b8+b4)>0,(b8-b4)/(b8+b4),np.nan)
        if indice in ('gndvi','ambos'): r['GNDVI']=np.where((b8+b3)>0,(b8-b3)/(b8+b3),np.nan)
    return r

def val_punto(img,lat,lon,olat,olon,plat,plon):
    rows,cols=img.shape; c=int((lon-olon)/plon); r=int((olat-lat)/plat)
    r0,r1=max(0,r-1),min(rows,r+2); c0,c1=max(0,c-1),min(cols,c+2)
    if r0>=r1 or c0>=c1: return None
    v=float(np.nanmean(img[r0:r1,c0:c1])); return v if np.isfinite(v) else None

def val_lote(img,poly,olat,olon,plat,plon):
    if isinstance(poly, dict) and 'multipart' in poly:
        all_vals = []
        for p in poly['multipart']:
            v, n = val_lote(img, p, olat, olon, plat, plon)
            if v is not None:
                all_vals.extend([v]*n)
        return (float(sum(all_vals)/len(all_vals)), len(all_vals)) if all_vals else (None, 0)
    rows,cols=img.shape
    lons=[p[0] for p in poly]; lats=[p[1] for p in poly]
    c0=max(0,int((min(lons)-olon)/plon)-2); c1=min(cols-1,int((max(lons)-olon)/plon)+2)
    r0=max(0,int((olat-max(lats))/plat)-2); r1=min(rows-1,int((olat-min(lats))/plat)+2)
    vals=[]
    for r in range(r0,r1+1):
        for c in range(c0,c1+1):
            if pip(olon+c*plon, olat-r*plat, poly):
                v=img[r,c]
                if np.isfinite(v) and v>0: vals.append(float(v))
    return (float(np.mean(vals)),len(vals)) if vals else (None,0)

# MAPA
def gen_mapa(img,poly,puntos,olat,olon,plat,plon,titulo,idx_nom,unidad):
    rows,cols=img.shape; mg=0.003
    # Normalizar: siempre trabajar con lista de partes
    partes=poly['multipart'] if isinstance(poly,dict) and 'multipart' in poly else [poly]
    all_lons=[p[0] for pp in partes for p in pp]
    all_lats=[p[1] for pp in partes for p in pp]
    c0=max(0,int((min(all_lons)-mg-olon)/plon)); c1=min(cols-1,int((max(all_lons)+mg-olon)/plon)+1)
    r0=max(0,int((olat-max(all_lats)-mg)/plat)); r1=min(rows-1,int((olat-min(all_lats)+mg)/plat)+1)
    crop=img[r0:r1+1,c0:c1+1].copy(); h,w=crop.shape
    mask=np.zeros((h,w),dtype=bool)
    for r in range(h):
        for c in range(w):
            if pip(olon+(c0+c)*plon,olat-(r0+r)*plat,poly): mask[r,c]=True
    crop[~mask]=np.nan
    ext=[olon+c0*plon,olon+c1*plon,olat-r1*plat,olat-r0*plat]
    fig,ax=plt.subplots(figsize=(7,6),facecolor='white')
    cm2=plt.cm.RdYlGn.copy(); cm2.set_bad('#e8e8e8')
    im=ax.imshow(crop,cmap=cm2,vmin=0.5,vmax=1.0,extent=ext,origin='upper',interpolation='nearest')
    # Cada parte se dibuja como poligono cerrado e independiente (sin linea entre partes)
    for pp in partes:
        ax.plot([p[0] for p in pp]+[pp[0][0]],[p[1] for p in pp]+[pp[0][1]],'k-',lw=1.5,zorder=5)
    clrs=['#1565C0','#E65100','#2E7D32','#6A1B9A','#AD1457','#00838F']
    # Calcular offsets inteligentes para evitar superposicion
    from itertools import product as iproduct
    posiciones_usadas = []
    offsets_candidatos = [(10,10),(10,-22),(-60,10),(-60,-22),(10,22),(-60,22),(30,-10),(-80,10)]

    for i,(amb,lat,lon,rinde) in enumerate(puntos):
        cp=clrs[i%len(clrs)]
        ax.scatter(lon,lat,s=130,c=cp,zorder=10,edgecolors='white',linewidth=1.5)
        # Elegir offset que no se superponga con etiquetas anteriores
        mejor_off = offsets_candidatos[i % len(offsets_candidatos)]
        for ox,oy in offsets_candidatos:
            solapado = False
            for px,py in posiciones_usadas:
                if abs((lon+ox*0.0001)-px) < 0.003 and abs((lat+oy*0.0001)-py) < 0.002:
                    solapado = True; break
            if not solapado:
                mejor_off = (ox,oy); break
        posiciones_usadas.append((lon+mejor_off[0]*0.0001, lat+mejor_off[1]*0.0001))
        rinde_fmt = fmt(rinde, 1) if isinstance(rinde, (int,float)) else str(rinde)
        ax.annotate(f"{amb}\n{rinde_fmt} {unidad}",xy=(lon,lat),
                    xytext=mejor_off,textcoords='offset points',
                    fontsize=7.5,fontweight='bold',color=cp,
                    arrowprops=dict(arrowstyle='-',color=cp,lw=0.8),
                    path_effects=[pe.withStroke(linewidth=2,foreground='white')],zorder=11)
    cb=plt.colorbar(im,ax=ax,fraction=0.03,pad=0.02); cb.set_label(idx_nom,fontsize=9); cb.ax.tick_params(labelsize=8)
    ax.set_xlabel('Longitud',fontsize=8); ax.set_ylabel('Latitud',fontsize=8); ax.tick_params(labelsize=7)
    ax.set_title(titulo,fontsize=10,fontweight='bold',pad=8); plt.tight_layout()
    buf=io.BytesIO(); plt.savefig(buf,dpi=150,bbox_inches='tight',facecolor='white'); plt.close(); buf.seek(0)
    return buf

# PDF
def gen_pdf(lotes,config,path):
    from datetime import datetime
    import base64 as _b64
    # Preparar logo para el PDF
    logo_buf = io.BytesIO(_b64.b64decode(LOGO_B64))
    LOGO_W = 2.5*cm
    LOGO_H = 1.9*cm

    GD=colors.HexColor('#1B5E20'); GM=colors.HexColor('#2E7D32'); GC=colors.HexColor('#C8E6C9')
    GL=colors.HexColor('#F5F5F5'); LN=colors.HexColor('#CFD8DC'); OR=colors.HexColor('#E65100')
    WH=colors.white; GT=colors.HexColor('#37474F')
    def s(n,**k):
        d=dict(fontName='Helvetica',fontSize=10,textColor=GT,leading=14,spaceAfter=0,spaceBefore=0); d.update(k); return ParagraphStyle(n,**d)
    s_tit=s('t',fontSize=20,textColor=GD,fontName='Helvetica-Bold',alignment=TA_CENTER,leading=26,spaceAfter=8)
    s_sub=s('sb',fontSize=11,alignment=TA_CENTER,leading=15,spaceAfter=4)
    s_fch=s('f',fontSize=9,textColor=colors.HexColor('#78909C'),alignment=TA_CENTER,leading=13)
    s_sec=s('sc',fontSize=10,fontName='Helvetica-Bold',textColor=GM,leading=14,spaceAfter=5,spaceBefore=8)
    s_not=s('n',fontSize=8,fontName='Helvetica-Oblique',textColor=colors.HexColor('#78909C'),leading=11,spaceAfter=3)
    s_fot=s('fo',fontSize=7.5,textColor=colors.HexColor('#90A4AE'),alignment=TA_CENTER,leading=11)
    s_rl=s('rl',fontSize=10,fontName='Helvetica-Bold',textColor=WH,alignment=TA_CENTER,leading=14)
    s_rv=s('rv',fontSize=16,fontName='Helvetica-Bold',textColor=WH,alignment=TA_CENTER,leading=22)
    s_pl=s('pl',fontSize=10,fontName='Helvetica-Bold',textColor=WH,alignment=TA_CENTER,leading=14)
    s_pv=s('pv',fontSize=13,fontName='Helvetica-Bold',textColor=WH,alignment=TA_CENTER,leading=18)
    HDR=[('BACKGROUND',(0,0),(-1,0),GM),('TEXTCOLOR',(0,0),(-1,0),WH),('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
         ('FONTSIZE',(0,0),(-1,-1),9),('ALIGN',(0,0),(-1,-1),'CENTER'),('VALIGN',(0,0),(-1,-1),'MIDDLE'),
         ('FONTNAME',(0,1),(-1,-1),'Helvetica'),('ROWBACKGROUNDS',(0,1),(-1,-1),[GL,WH]),('GRID',(0,0),(-1,-1),0.4,LN),
         ('TOPPADDING',(0,0),(-1,-1),6),('BOTTOMPADDING',(0,0),(-1,-1),6),('LEFTPADDING',(0,0),(-1,-1),8),('RIGHTPADDING',(0,0),(-1,-1),8)]
    def tbl(data,cw,ex=None):
        t=Table(data,colWidths=cw,repeatRows=1); t.setStyle(TableStyle(list(HDR)+(ex or []))); return t
    indice=config['indice'].upper(); unidad=config['unidad']; now=datetime.now().strftime('%d/%m/%Y')
    doc=SimpleDocTemplate(path,pagesize=A4,leftMargin=2*cm,rightMargin=2*cm,topMargin=1.8*cm,bottomMargin=1.8*cm)
    W=A4[0]-4*cm; story=[]
    # Encabezado con logo
    enc = Table([
        [Paragraph(config['titulo'], s_tit), RLImage(logo_buf, width=LOGO_W, height=LOGO_H)]
    ], colWidths=[W - LOGO_W - 0.3*cm, LOGO_W + 0.3*cm])
    enc.setStyle(TableStyle([
        ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
        ('ALIGN',(1,0),(1,0),'RIGHT'),
        ('LEFTPADDING',(0,0),(-1,-1),0),
        ('RIGHTPADDING',(0,0),(-1,-1),0),
        ('TOPPADDING',(0,0),(-1,-1),0),
        ('BOTTOMPADDING',(0,0),(-1,-1),6),
        ('LINEBELOW',(0,0),(-1,0),2,GD),
    ]))
    story += [enc,
        Paragraph(f"Destino: {config['destino']} &nbsp;&middot;&nbsp; Analisis {indice} &nbsp;&middot;&nbsp; Sentinel-2",s_sub),
        Spacer(1,4),Paragraph(f"Imagen: {config['fecha']} &nbsp;&middot;&nbsp; Generado: {now}",s_fch),
        Spacer(1,12)]
    story.append(Paragraph('Resumen por Lote',s_sec))
    dr=[['Lote','Campo','Cultivo','Sup. (ha)',indice+' prom.','R²',f'Rinde ({unidad})']]
    for l in lotes:
        dr.append([l['nombre'],l['campo'],l['cultivo'],fmt(l['area_ha'], 1),f"{l['idx_prom']:.4f}",
                   f"{l['r2']:.3f}"+(' *' if l['r2']<0.5 else ''),fmt(l['rinde_est'], 2)])
    story.append(tbl(dr,[2.7*cm,3*cm,1.9*cm,2*cm,2.2*cm,1.3*cm,2.9*cm],
                     [('TEXTCOLOR',(6,1),(6,-1),GD),('FONTNAME',(6,1),(6,-1),'Helvetica-Bold')]))
    story.append(Spacer(1,5))
    if any(l['r2']<0.5 for l in lotes): story.append(Paragraph('* R² bajo: correlacion debil.',s_not))
    story.append(Spacer(1,16))
    story+=[HRFlowable(width=W,thickness=1,color=LN,spaceAfter=12),Paragraph('Produccion Total Estimada',s_sec)]
    at=sum(l['area_ha'] for l in lotes); pt=sum(l['rinde_est']*l['area_ha']/1000.0 for l in lotes); rp=sum(l['rinde_est']*l['area_ha'] for l in lotes)/at if at else 0
    dp=[['Lote','Superficie (ha)',f'Rinde ({unidad})','Produccion (tn)']]
    for l in lotes: dp.append([l['nombre'],fmt(l['area_ha'], 1),fmt(l['rinde_est'], 2),fmt(l['rinde_est']*l['area_ha'], 1)])
    dp.append(['TOTAL / PROM. POND.',fmt(at, 1),fmt(rp, 2),fmt(pt, 1)])
    story.append(tbl(dp,[4.5*cm,3.5*cm,3.5*cm,4.5*cm],
                     [('BACKGROUND',(0,-1),(-1,-1),GD),('TEXTCOLOR',(0,-1),(-1,-1),WH),('FONTNAME',(0,-1),(-1,-1),'Helvetica-Bold')]))
    story.append(Spacer(1,16))
    cx=Table([[Paragraph('PRODUCCION TOTAL ESTIMADA',s_pl)],
              [Paragraph(f"{fmt(pt,1)} tn &nbsp;|&nbsp; Rinde pond.: {fmt(rp,2)} {unidad} &nbsp;|&nbsp; Sup.: {fmt(at,1)} ha",s_pv)]],colWidths=[W])
    cx.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1),GD),('TOPPADDING',(0,0),(0,0),9),('BOTTOMPADDING',(0,0),(0,0),4),
                             ('TOPPADDING',(0,1),(0,1),4),('BOTTOMPADDING',(0,1),(0,1),10),
                             ('LEFTPADDING',(0,0),(-1,-1),12),('RIGHTPADDING',(0,0),(-1,-1),12),('LINEABOVE',(0,0),(-1,0),3,OR)]))
    story.append(cx)
    IW=W; IH=W*0.60
    for l in lotes:
        story.append(PageBreak())
        ht=Table([[Paragraph(f"<b>{l['nombre']}</b>",s('lh',fontSize=14,fontName='Helvetica-Bold',textColor=WH,alignment=TA_LEFT,leading=18)),
                   Paragraph(f"Campo: {l['campo']} &nbsp;|&nbsp; Cultivo: {l['cultivo']} &nbsp;|&nbsp; Sup.: {l['area_ha']:.2f} ha",
                              s('lhi',fontSize=9,textColor=GC,alignment=TA_LEFT,leading=13))]],colWidths=[4*cm,W-4*cm])
        ht.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1),GD),('VALIGN',(0,0),(-1,-1),'MIDDLE'),
                                 ('TOPPADDING',(0,0),(-1,-1),10),('BOTTOMPADDING',(0,0),(-1,-1),10),
                                 ('LEFTPADDING',(0,0),(-1,-1),12),('RIGHTPADDING',(0,0),(-1,-1),8),('LINEBELOW',(0,0),(-1,-1),3,OR)]))
        story+=[ht,Spacer(1,10),Paragraph(f'Imagen {indice} y Puntos de Muestreo',s_sec)]
        l['mapa_buf'].seek(0); story.append(RLImage(l['mapa_buf'],width=IW,height=IH)); story.append(Spacer(1,10))
        pd2=[['Ambiente',indice,f'Rinde ({unidad})']]
        for p in l['puntos_datos']: pd2.append([p['amb'],f"{p['idx_val']:.4f}",fmt(p['rinde'], 2)])
        story.append(tbl(pd2,[6*cm,5*cm,6*cm])); story.append(Spacer(1,12))
        story+=[HRFlowable(width=W,thickness=1,color=LN,spaceAfter=8),Paragraph(f'Modelo de Correlacion {indice} - Rendimiento',s_sec)]
        dm=[['Ecuacion','R²',f'{indice} prom.'],[l['ecuacion'],f"{l['r2']:.3f}"+(' *' if l['r2']<0.5 else ''),f"{l['idx_prom']:.4f}"]]
        story.append(tbl(dm,[8*cm,2.5*cm,5.5*cm])); story.append(Spacer(1,6))
        if l['r2']<0.5: story.append(Paragraph('* R² bajo: interpretar con precaucion.',s_not))
        story.append(Spacer(1,6))
        rt=Table([[Paragraph('RENDIMIENTO ESTIMADO',s_rl)],
                  [Paragraph(f"{fmt(l['rinde_est'],2)} {unidad} &nbsp;|&nbsp; {fmt(l['rinde_est']*l['area_ha']/1000.0,1)} tn totales",s_rv)]],colWidths=[W])
        rt.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1),GM),('TOPPADDING',(0,0),(0,0),8),('BOTTOMPADDING',(0,0),(0,0),4),
                                 ('TOPPADDING',(0,1),(0,1),4),('BOTTOMPADDING',(0,1),(0,1),10),
                                 ('LEFTPADDING',(0,0),(-1,-1),12),('RIGHTPADDING',(0,0),(-1,-1),12),('LINEABOVE',(0,0),(-1,0),3,OR)]))
        story+=[rt,Spacer(1,12),HRFlowable(width=W,thickness=0.5,color=LN,spaceAfter=5),
                Paragraph(f'Sentinel-2 L2A (Copernicus) | {indice} | Regresion lineal | {now}',s_fot)]
    # ── Resumen general final ─────────────────────────────────
    story.append(PageBreak())
    enc2_buf = io.BytesIO(_b64.b64decode(LOGO_B64))
    enc2 = Table([
        [Paragraph('Resumen General — Campo '+lotes[0]['campo'] if lotes else 'Resumen General', 
                   s('rg', fontSize=16, fontName='Helvetica-Bold', textColor=GD, leading=20)),
         RLImage(enc2_buf, width=LOGO_W, height=LOGO_H)]
    ], colWidths=[W - LOGO_W - 0.3*cm, LOGO_W + 0.3*cm])
    enc2.setStyle(TableStyle([
        ('VALIGN',(0,0),(-1,-1),'MIDDLE'),('ALIGN',(1,0),(1,0),'RIGHT'),
        ('LEFTPADDING',(0,0),(-1,-1),0),('RIGHTPADDING',(0,0),(-1,-1),0),
        ('TOPPADDING',(0,0),(-1,-1),0),('BOTTOMPADDING',(0,0),(-1,-1),6),
        ('LINEBELOW',(0,0),(-1,0),2,GD),
    ]))
    story += [enc2, Spacer(1,6),
              Paragraph(f'Todos los lotes analizados &nbsp;&middot;&nbsp; Generado: {now}', s_fch),
              Spacer(1,14)]

    # Agrupar por variedad para el resumen
    from collections import OrderedDict
    grupos_var = OrderedDict()
    for l in lotes:
        k = (l.get('variedad',''), l.get('cultivo',''))
        if k not in grupos_var: grupos_var[k] = []
        grupos_var[k].append(l)

    for (var, cult), lotes_g in grupos_var.items():
        label = cult + (f' — Var: {var}' if var else '')
        story.append(Paragraph(label, s_sec))
        dr2 = [['Lote','Campo','Sup.(ha)',indice+' prom.','R²',f'Rinde ({unidad})','Produccion (tn)']]
        for l in lotes_g:
            dr2.append([l['nombre'], l['campo'], fmt(l['area_ha'],1), f"{l['idx_prom']:.4f}",
                        f"{l['r2']:.3f}"+(' *' if l['r2']<0.5 else ''),
                        fmt(l['rinde_est'],2), fmt(l['rinde_est']*l['area_ha']/1000.0,1)])
        at_g = sum(x['area_ha'] for x in lotes_g)
        pt_g = sum(x['rinde_est']*x['area_ha']/1000.0 for x in lotes_g)
        rp_g = pt_g/at_g if at_g else 0
        dr2.append(['SUBTOTAL','','',f"{fmt(at_g,1)} ha",'',fmt(rp_g,2)+f' {unidad}',fmt(pt_g,1)+' tn'])
        t2=Table(dr2,colWidths=[2.5*cm,2.5*cm,1.8*cm,2*cm,1.2*cm,2.5*cm,3*cm],repeatRows=1)
        t2.setStyle(TableStyle(list(HDR)+[
            ('TEXTCOLOR',(5,1),(5,-1),GD),('FONTNAME',(5,1),(5,-1),'Helvetica-Bold'),
            ('BACKGROUND',(0,-1),(-1,-1),GM),('TEXTCOLOR',(0,-1),(-1,-1),WH),
            ('FONTNAME',(0,-1),(-1,-1),'Helvetica-Bold'),
        ]))
        story += [t2, Spacer(1,12)]

    # Caja total campo completo
    at_tot = sum(l['area_ha'] for l in lotes)
    pt_tot = sum(l['rinde_est']*l['area_ha']/1000.0 for l in lotes)
    rp_tot = sum(l['rinde_est']*l['area_ha'] for l in lotes)/at_tot if at_tot else 0
    campo_nom = lotes[0]['campo'] if lotes else ''
    ct=Table([[Paragraph(f'RESUMEN CAMPO {campo_nom.upper()}',s_pl)],
              [Paragraph(f"{fmt(pt_tot,1)} tn &nbsp;|&nbsp; Rinde pond.: {fmt(rp_tot,2)} {unidad} &nbsp;|&nbsp; Sup. total activa: {fmt(at_tot,1)} ha",s_pv)]],
             colWidths=[W])
    ct.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,-1),GD),
        ('TOPPADDING',(0,0),(0,0),9),('BOTTOMPADDING',(0,0),(0,0),4),
        ('TOPPADDING',(0,1),(0,1),4),('BOTTOMPADDING',(0,1),(0,1),10),
        ('LEFTPADDING',(0,0),(-1,-1),12),('RIGHTPADDING',(0,0),(-1,-1),12),
        ('LINEABOVE',(0,0),(-1,0),3,OR)]))
    story += [ct, Spacer(1,10),
              Paragraph('Superficie activa = píxeles con GNDVI ≥ 0,5 dentro del perímetro de cada lote.', s_not)]

    doc.build(story)

# ═══════════════════════════════════════════════
# POWERPOINT (python-pptx)
# ═══════════════════════════════════════════════

def gen_pptx_python(lotes, config, path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import base64 as _b

    VERDE_OSC = RGBColor(0x1B, 0x5E, 0x20)
    VERDE_MED = RGBColor(0x2E, 0x7D, 0x32)
    VERDE_CLA = RGBColor(0xC8, 0xE6, 0xC9)
    NARANJA   = RGBColor(0xE6, 0x51, 0x00)
    GRIS_TX   = RGBColor(0x37, 0x47, 0x4F)
    GRIS_CLR  = RGBColor(0xF5, 0xF5, 0xF5)
    BLANCO    = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = 13.33
    H = 7.5

    logo_bytes = _b.b64decode(LOGO_B64)
    unidad = config.get('unidad', 'kg/ha')

    def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
        from pptx.util import Inches
        shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb
        if line_rgb:
            shp.line.color.rgb = line_rgb
        else:
            shp.line.fill.background()
        return shp

    def add_txt(slide, text, x, y, w, h, size=12, bold=False, color=None, align='left', italic=False):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
        run.font.color.rgb = color or GRIS_TX
        return txb

    def add_header(slide, titulo, subtitulo=''):
        add_rect(slide, 0, 0, W, 0.7, VERDE_OSC)
        add_rect(slide, 0, 0.7, W, 0.06, NARANJA)
        add_txt(slide, titulo, 0.3, 0.05, W-1.8, 0.6, size=20, bold=True, color=BLANCO, align='left')
        logo_buf2 = io.BytesIO(logo_bytes)
        slide.shapes.add_picture(logo_buf2, Inches(W-1.5), Inches(0.05), Inches(1.3), Inches(0.6))
        if subtitulo:
            add_txt(slide, subtitulo, 0.3, 0.76, W-0.6, 0.28, size=9, color=GRIS_TX, italic=True)

    def add_footer(slide):
        add_txt(slide, f"Sentinel-2 L2A (Copernicus) | GNDVI=(B8-B3)/(B8+B3) | Regresion lineal | {config.get('fecha','')}",
                0.3, H-0.28, W-0.6, 0.25, size=7, color=RGBColor(0x90,0xA4,0xAE), italic=True)

    def tbl_cell(table, row, col, text, bold=False, bg=None, color=None, align='center'):
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        cell = table.cell(row, col)
        cell.text = str(text)
        p = cell.text_frame.paragraphs[0]
        p.alignment = {'center':PP_ALIGN.CENTER,'left':PP_ALIGN.LEFT,'right':PP_ALIGN.RIGHT}.get(align,PP_ALIGN.CENTER)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = bold; run.font.size = Pt(9)
        run.font.color.rgb = color or GRIS_TX
        if bg:
            cell.fill.solid(); cell.fill.fore_color.rgb = bg

    blank_layout = prs.slide_layouts[6]

    # ── Portada ───────────────────────────────────────────
    s0 = prs.slides.add_slide(blank_layout)
    s0.background.fill.solid(); s0.background.fill.fore_color.rgb = VERDE_OSC
    logo_buf0 = io.BytesIO(logo_bytes)
    s0.shapes.add_picture(logo_buf0, Inches(5.67), Inches(0.4), Inches(2.0), Inches(1.5))
    add_txt(s0, config.get('titulo','Estimacion de Rendimiento'),
            0.5, 2.1, 12.3, 1.2, size=30, bold=True, color=BLANCO, align='center')
    add_rect(s0, 4.2, 3.45, 4.9, 0.07, NARANJA)
    at = sum(l['area_ha'] for l in lotes)
    pt = sum(l['rinde_est']*l['area_ha'] for l in lotes)
    rp = pt/at if at else 0
    add_txt(s0, f"Destino: {config.get('destino','')}  ·  GNDVI  ·  Sentinel-2  ·  {config.get('fecha','')}",
            0.5, 3.6, 12.3, 0.4, size=13, color=VERDE_CLA, align='center')
    add_rect(s0, 2.0, 4.2, 9.33, 0.9, RGBColor(0x15,0x5E,0x20))
    add_txt(s0, f"{fmt(pt,1)} tn  ·  Rinde pond.: {fmt(rp,2)} {unidad}  ·  Superficie: {fmt(at,1)} ha",
            2.0, 4.2, 9.33, 0.9, size=14, bold=True, color=BLANCO, align='center')

    # ── Resumen ───────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    add_header(s1, 'Resumen por Lote', config.get('titulo',''))
    cols = ['Lote','Campo','Variedad','Sup.(ha)','GNDVI prom.','R²',f'Rinde ({unidad})','Prod.(tn)']
    n_cols = len(cols); n_rows = len(lotes)+2
    tbl = s1.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(1.1), Inches(W-0.6), Inches(H-1.6)).table
    col_widths = [1.8,1.8,1.4,0.9,1.1,0.7,1.3,1.2]
    for i,cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)
    for ci,ch in enumerate(cols):
        tbl_cell(tbl,0,ci,ch,bold=True,bg=VERDE_MED,color=BLANCO)
    for ri,l in enumerate(lotes):
        bg = GRIS_CLR if ri%2==0 else BLANCO
        vals = [l['nombre'],l['campo'],l.get('variedad','-'),fmt(l['area_ha'],1),
                f"{l['idx_prom']:.4f}",f"{l['r2']:.3f}"+(' *' if l['r2']<0.5 else ''),
                fmt(l['rinde_est'],2), fmt(l['rinde_est']*l['area_ha']/1000.0,1)]
        for ci,v in enumerate(vals):
            tbl_cell(tbl,ri+1,ci,v,bg=bg,color=VERDE_OSC if ci==6 else GRIS_TX)
    tot_vals=['TOTAL','','','',fmt(at,1),'',fmt(rp,2),fmt(pt,1)+' tn']
    for ci,v in enumerate(tot_vals):
        tbl_cell(tbl,n_rows-1,ci,v,bold=True,bg=VERDE_OSC,color=BLANCO)
    add_footer(s1)

    # ── Una slide por lote ────────────────────────────────
    for l in lotes:
        sl = prs.slides.add_slide(blank_layout)
        var_str = f" · Var: {l['variedad']}" if l.get('variedad') else ''
        add_header(sl, l['nombre'],
                   f"Campo: {l['campo']}{var_str}  ·  Cultivo: {l['cultivo']}  ·  Sup.: {fmt(l['area_ha'],1)} ha")

        # Mapa
        if l.get('mapa_buf'):
            try:
                l['mapa_buf'].seek(0)
                sl.shapes.add_picture(l['mapa_buf'], Inches(0.3), Inches(1.1), Inches(7.5), Inches(5.0))
            except: pass

        # Panel derecho - solo puntos dentro del lote (ya filtrados en procesar)
        px = 8.1
        add_rect(sl, px, 1.1, 4.9, 0.32, VERDE_MED)
        add_txt(sl, 'Puntos de Muestreo', px, 1.1, 4.9, 0.32, size=9, bold=True, color=BLANCO, align='center')

        pts = l.get('puntos_datos', [])
        if pts:
            pt_tbl = sl.shapes.add_table(len(pts)+1, 3, Inches(px), Inches(1.44),
                                          Inches(4.9), Inches(0.35+len(pts)*0.3)).table
            pt_tbl.columns[0].width = Inches(2.0)
            pt_tbl.columns[1].width = Inches(1.3)
            pt_tbl.columns[2].width = Inches(1.6)
            for ci,ch in enumerate(['Ambiente','GNDVI',f'Rinde ({unidad})']):
                tbl_cell(pt_tbl,0,ci,ch,bold=True,bg=VERDE_MED,color=BLANCO)
            for ri,p in enumerate(pts):
                bg = GRIS_CLR if ri%2==0 else BLANCO
                tbl_cell(pt_tbl,ri+1,0,p['amb'],bg=bg,color=GRIS_TX,align='left')
                tbl_cell(pt_tbl,ri+1,1,f"{p['idx_val']:.4f}",bg=bg,color=GRIS_TX)
                tbl_cell(pt_tbl,ri+1,2,fmt(p['rinde'],2),bg=bg,color=GRIS_TX)

        y_mod = 1.46 + 0.35 + len(pts)*0.3 + 0.15
        add_rect(sl, px, y_mod, 4.9, 0.3, VERDE_MED)
        add_txt(sl, 'Modelo de Correlacion', px, y_mod, 4.9, 0.3, size=9, bold=True, color=BLANCO, align='center')
        add_txt(sl, l.get('ecuacion',''), px+0.1, y_mod+0.32, 3.2, 0.28, size=8, color=GRIS_TX)
        add_txt(sl, f"R²: {l['r2']:.3f}", px+3.3, y_mod+0.32, 1.5, 0.28, size=8, color=GRIS_TX)
        add_txt(sl, f"GNDVI prom.: {l['idx_prom']:.4f}", px+0.1, y_mod+0.62, 4.7, 0.28, size=8, color=GRIS_TX)

        # Caja rendimiento
        y_rinde = H - 1.5
        add_rect(sl, px, y_rinde, 4.9, 1.3, VERDE_MED)
        add_txt(sl, 'RENDIMIENTO ESTIMADO', px, y_rinde+0.05, 4.9, 0.28,
                size=9, bold=True, color=VERDE_CLA, align='center')
        add_txt(sl, f"{fmt(l['rinde_est'],2)} {unidad}", px, y_rinde+0.32, 4.9, 0.55,
                size=26, bold=True, color=BLANCO, align='center')
        add_txt(sl, f"Produccion: {fmt(l['rinde_est']*l['area_ha']/1000.0,1)} tn  ·  Sup.: {fmt(l['area_ha'],1)} ha",
                px, y_rinde+0.9, 4.9, 0.3, size=9, color=VERDE_CLA, align='center')
        add_footer(sl)

    prs.save(path)

# HTTP
class Handler(http.server.BaseHTTPRequestHandler):
    def log_message(self,fmt,*a): pass
    def cors(self):
        self.send_header('Access-Control-Allow-Origin','*')
        self.send_header('Access-Control-Allow-Methods','POST, GET, OPTIONS')
        self.send_header('Access-Control-Allow-Headers','Content-Type')
    def do_OPTIONS(self): self.send_response(200); self.cors(); self.end_headers()
    def do_GET(self):
        if self.path=='/status':
            self.send_response(200); self.send_header('Content-Type','application/json'); self.cors(); self.end_headers()
            r={'ok':LIBS_OK}
            if not LIBS_OK: r['error']=LIBS_ERR
            self.wfile.write(json.dumps(r).encode())
        elif self.path in ('/', '/app', '/app.html', '/index.html'):
            html=_get_html(); self.send_response(200); self.send_header('Content-Type','text/html; charset=utf-8'); self.cors(); self.end_headers()
            self.wfile.write(html.encode('utf-8'))

    def do_POST(self):
        length=int(self.headers.get('Content-Length',0)); body=self.rfile.read(length)
        try:
            result=self.procesar(json.loads(body))
            self.send_response(200); self.send_header('Content-Type','application/json'); self.cors(); self.end_headers()
            self.wfile.write(json.dumps(result).encode())
        except Exception as e:
            import traceback
            tb = traceback.format_exc()
            print("\n=== ERROR ===")
            print(tb)
            print("=============\n")
            self.send_response(500); self.send_header('Content-Type','application/json'); self.cors(); self.end_headers()
            self.wfile.write(json.dumps({'error':str(e),'detalle':tb}).encode())
    def procesar(self,payload):
        if not LIBS_OK: raise Exception(f'Falta libreria: {LIBS_ERR}. Ejecuta ARRANCAR.bat.')
        config=payload['config']; grupos=payload['grupos']
        records=leer_dbf(base64.b64decode(payload['dbf']))
        polys=leer_shp(base64.b64decode(payload['shp']))
        cl=col_lote(records)
        b3,b4,b8,geo=leer_zip(base64.b64decode(payload['zip']))
        print(f"DEBUG leer_zip retorno: b3={b3 is not None}, b4={b4 is not None}, b8={b8 is not None}, geo={geo}")
        if geo is None:
            raise Exception("geo es None despues de leer_zip - el ZIP no pudo ser leido correctamente")
        olon,olat,plon,plat=geo
        print(f"DEBUG geo: olon={olon:.6f}, olat={olat:.6f}, plon={plon:.8f}, plat={plat:.8f}")
        print(f"DEBUG imagen shape: {b8.shape}, bbox lat=[{olat-b8.shape[0]*plat:.4f}, {olat:.4f}], lon=[{olon:.4f}, {olon+b8.shape[1]*plon:.4f}]")
        indices=calc_idx(b3,b4,b8,config['indice'])
        idx_nom='GNDVI' if config['indice']=='gndvi' else 'NDVI' if config['indice']=='ndvi' else 'GNDVI'
        img=indices.get('GNDVI',indices.get('NDVI'))
        print(f"DEBUG indices keys: {list(indices.keys())}, img is None: {img is None}")
        if img is None:
            raise Exception(f"No se pudo calcular el indice {config['indice']}. Verificá que las bandas correctas esten en el ZIP.")
        lotes_res=[]
        for grupo in grupos:
            pts=[]
            for pt in grupo['puntos']:
                try:
                    lat_v = float(pt['lat'])
                    lon_v = float(pt['lon'])
                    rinde_v = float(pt['rinde'])
                except (TypeError, ValueError):
                    print(f"Punto ignorado (valor invalido): {pt}")
                    continue
                import math
                if math.isnan(lat_v) or math.isnan(lon_v) or math.isnan(rinde_v):
                    print(f"Punto ignorado (NaN): amb={pt['amb']}")
                    continue
                print(f"Punto: {pt['amb']} lat={lat_v} lon={lon_v} rinde={rinde_v}")
                v=val_punto(img, lat_v, lon_v, olat, olon, plat, plon)
                print(f"  -> val_punto={v}")
                if v is not None:
                    pts.append({'amb':pt['amb'],'lat':lat_v,'lon':lon_v,'rinde':rinde_v,'idx_val':v})
            if len(pts)<2: raise Exception(f'El grupo "{grupo["nombre"]}" necesita al menos 2 puntos validos en la imagen.')
            xg=np.array([p['idx_val'] for p in pts]); yg=np.array([p['rinde'] for p in pts])
            sl,it,r,_,_=stats.linregress(xg,yg); r2=r**2
            for lote in grupo['lotes']:
                poly=buscar_poly(lote['nombre'],records,polys,cl)
                if poly is None:
                    disp=list(set(rec.get(cl,'') for rec in records if rec.get(cl,'')))
                    raise Exception(f'Lote "{lote["nombre"]}" no encontrado.\nColumna: "{cl}".\nDisponibles: {disp[:25]}')
                ha_manual = float(lote.get('sup_ha', 0) or 0)
                ha = ha_manual if ha_manual > 0 else area_ha(poly)
                prom,_=val_lote(img,poly,olat,olon,plat,plon)
                if prom is None: raise Exception(f'Sin pixeles validos en lote "{lote["nombre"]}".')
                rinde_kgha = sl*prom+it
                rinde = convertir_rinde(rinde_kgha, config['unidad'])
                ec=f"y = {sl:.2f} x {idx_nom} + ({it:.2f})"  
                # Solo mostrar en el mapa los puntos que caen dentro de este lote
                pts_en_lote = [(p['amb'],p['lat'],p['lon'],convertir_rinde(p['rinde'],config['unidad']))
                               for p in pts if pip(p['lon'], p['lat'], poly)]
                mb=gen_mapa(img,poly,pts_en_lote,olat,olon,plat,plon,
                            f"{idx_nom} - {lote['nombre']}",idx_nom,config['unidad'])
                lotes_res.append({'nombre':lote['nombre'],'campo':lote['campo'],'cultivo':lote['cultivo'],
                                  'variedad':grupo.get('variedad',''),'grupo_nombre':grupo['nombre'],
                                  'area_ha':ha,'idx_prom':prom,'r2':r2,'rinde_est':rinde,'ecuacion':ec,
                                  'puntos_datos':[dict(p, rinde=convertir_rinde(p['rinde'],config['unidad'])) for p in pts],'mapa_buf':mb})
        with tempfile.NamedTemporaryFile(suffix='.pdf',delete=False) as tmp: pdf_path=tmp.name
        gen_pdf(lotes_res,config,pdf_path)
        with open(pdf_path,'rb') as f: b64=base64.b64encode(f.read()).decode()
        os.unlink(pdf_path)
        # ── Generar PPTX con python-pptx ─────────────────────
        import base64 as _b64_2
        pptx_b64 = ''
        try:
            pptx_tmp = pdf_path.replace('.pdf', '.pptx')
            gen_pptx_python(lotes_res, config, pptx_tmp)
            with open(pptx_tmp, 'rb') as f2:
                pptx_b64 = _b64_2.b64encode(f2.read()).decode()
            os.unlink(pptx_tmp)
        except Exception as epptx:
            print(f"PPTX error: {epptx}")

        return {'ok':True,'pdf':b64,'pptx':pptx_b64}

if __name__=='__main__':
    print(f'Servidor iniciando en puerto {PORT}')
    socketserver.TCPServer.allow_reuse_address = True
    with socketserver.TCPServer(('0.0.0.0', PORT), Handler) as srv:
        srv.serve_forever()
